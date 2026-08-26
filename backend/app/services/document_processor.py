"""Extract normalized text from development/testing artifacts.

Document Intelligence deliberately supports more than BRDs: requirements,
architecture, API contracts, spreadsheets, test artifacts and exports all feed
the same QA knowledge model. Binary visual-only assets remain in the Upload
Repository and can be handled by a future multimodal extractor.
"""
from __future__ import annotations

import csv
import io
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable, Dict

import docx
import openpyxl
import yaml
from pptx import Presentation
from pypdf import PdfReader
import xlrd


class UnsupportedDocumentTypeError(ValueError):
    pass


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    document = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(paragraphs)


def _extract_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_json(data: bytes) -> str:
    payload = json.loads(data)
    if isinstance(payload, dict) and "issues" in payload:
        issues = payload.get("issues") or []
        lines = []
        for issue in issues:
            fields = issue.get("fields", issue) if isinstance(issue, dict) else {}
            key = issue.get("key", fields.get("key", "")) if isinstance(issue, dict) else ""
            summary = fields.get("summary", "")
            description = fields.get("description", "")
            acceptance_criteria = fields.get("customfield_acceptance_criteria", "")
            labels = ", ".join(fields.get("labels", []) or []) if isinstance(fields.get("labels", []), list) else str(fields.get("labels", ""))
            comments = fields.get("comment", {}).get("comments", []) if isinstance(fields.get("comment"), dict) else []
            comment_text = "\n".join(str(c.get("body", "")) for c in comments if isinstance(c, dict))
            lines.append(
                f"[{key}] {summary}\nDescription: {description}\n"
                f"Acceptance Criteria: {acceptance_criteria}\nLabels: {labels}\n"
                f"Comments: {comment_text}\n"
            )
        return "\n---\n".join(lines)
    # OpenAPI, Postman collections, configuration exports and generic JSON
    # retain their keys so the AI can reason about contracts and mappings.
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _extract_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(str(value) for value in row) for row in reader)


def _extract_xlsx(data: bytes) -> str:
    workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"[SHEET: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_xls(data: bytes) -> str:
    workbook = xlrd.open_workbook(file_contents=data)
    lines: list[str] = []
    for sheet in workbook.sheets():
        lines.append(f"[SHEET: {sheet.name}]")
        for row_index in range(sheet.nrows):
            values = [str(sheet.cell_value(row_index, col)) for col in range(sheet.ncols)]
            if any(value.strip() for value in values):
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[SLIDE {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                lines.append(str(shape.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(lines)


def _extract_yaml(data: bytes) -> str:
    payload = yaml.safe_load(data.decode("utf-8", errors="replace"))
    return json.dumps(payload, ensure_ascii=False, indent=2, default=str)


def _extract_xml(data: bytes) -> str:
    root = ET.fromstring(data)
    lines: list[str] = []
    for element in root.iter():
        text = (element.text or "").strip()
        attrs = " ".join(f"{key}={value}" for key, value in element.attrib.items())
        if text or attrs:
            lines.append(f"{element.tag} {attrs}: {text}".strip())
    return "\n".join(lines)


def _extract_html(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    text = re.sub(r"<script\b[^>]*>.*?</script>", " ", text, flags=re.I | re.S)
    text = re.sub(r"<style\b[^>]*>.*?</style>", " ", text, flags=re.I | re.S)
    text = re.sub(r"<[^>]+>", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text


_EXTRACTORS: Dict[str, Callable[[bytes], str]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".txt": _extract_text,
    ".md": _extract_text,
    ".json": _extract_json,
    ".csv": _extract_csv,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xls,
    ".pptx": _extract_pptx,
    ".yaml": _extract_yaml,
    ".yml": _extract_yaml,
    ".xml": _extract_xml,
    ".html": _extract_html,
    ".htm": _extract_html,
}


def extract_text(filename: str, data: bytes) -> str:
    """Return normalized text extracted from one uploaded artifact."""
    extension = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(extension)
    if extractor is None:
        raise UnsupportedDocumentTypeError(
            f"Unsupported file extension '{extension}'. Supported: {sorted(_EXTRACTORS)}"
        )
    return extractor(data)
